"""슬라이드 4~10 생성 스크립트. 슬라이드 1~3은 유지."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SRC = r"d:\Unity\Projects\DungeonFront\Presentation\Presentation 9.pptx"
DST = SRC

FONT = "Malgun Gothic"
ACCENT = RGBColor(0x44, 0x72, 0xC4)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
MUTED = RGBColor(0x66, 0x66, 0x66)
GOLD = RGBColor(0xC9, 0xA0, 0x27)
REP = RGBColor(0x70, 0x30, 0xA0)


def delete_slide(prs, index: int) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    slide_id = slide_ids[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    prs.slides._sldIdLst.remove(slide_id)


def style_run(run, size=18, bold=False, color=DARK, name=FONT) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_title(slide, text: str) -> None:
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        style_run(run, size=32, bold=True, color=ACCENT)


def add_textbox(slide, left, top, width, height, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color)
    return box


def fill_table_cell(cell, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def shade_header_row(table, cols: int) -> None:
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_content_slide(prs, layout_idx=1):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def build_slide_4(prs):
    slide = add_content_slide(prs)
    set_title(slide, "하루 3단계 — 준비 · 생산 · 결산")

    add_textbox(
        slide,
        Inches(0.6), Inches(1.25), Inches(12.0), Inches(0.45),
        "공장은 하루 단위로 진행되며, 매일 아래 3단계를 반복합니다.",
        size=15, color=MUTED,
    )

    rows, cols = 4, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.8), Inches(12.0), Inches(3.6)).table
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(6.8)
    table.columns[2].width = Inches(3.6)

    headers = ["단계", "하는 일", "비고"]
    for i, h in enumerate(headers):
        fill_table_cell(table.cell(0, i), h, size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    shade_header_row(table, cols)

    data = [
        ("준비", "모험가 의뢰 수락 · 설비 재배치 · 레시피 설정", "탑다운 뷰"),
        ("생산", "현실 시간 5분간 공장 가동 · 수작업 · 고장 대응", "실시간 5분"),
        ("결산", "의뢰 납품 · 골드·명성 보상 수령", "별도 UI 화면"),
    ]
    for r, (phase, action, note) in enumerate(data, start=1):
        fill_table_cell(table.cell(r, 0), phase, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        fill_table_cell(table.cell(r, 1), action, size=14)
        fill_table_cell(table.cell(r, 2), note, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        Inches(2.0), Inches(5.65), Inches(9.3), Inches(0.5),
        "준비  →  생산 (5분)  →  결산  →  다음 날",
        size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )


def build_slide_5(prs):
    slide = add_content_slide(prs)
    set_title(slide, "생산 단계 — 현장 이벤트")

    add_textbox(
        slide,
        Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.55),
        "5분 동안 가만히 보기만 하는 게 아닙니다. 직접 뛰어다니며 공장을 살려야 합니다.",
        size=15, color=MUTED,
    )

    events = [
        ("🔧", "고장 수리", "기계가 고장나면\n망치로 두드려 수리"),
        ("🛠", "수동 가공", "자동화 전에는\n클릭으로 직접 완성"),
        ("📦", "물품 운반", "필요할 때\n직접 옮기기"),
    ]
    w = Inches(3.6)
    gap = Inches(0.45)
    start_x = Inches(0.85)
    top = Inches(2.0)
    h = Inches(4.2)

    for i, (icon, title, desc) in enumerate(events):
        left = start_x + i * (w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFA)
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1.5)

        add_textbox(slide, left, top + Inches(0.35), w, Inches(0.7), icon, size=40, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.2), w - Inches(0.4), Inches(0.5), title, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.85), w - Inches(0.4), Inches(1.8), desc, size=14, color=DARK, align=PP_ALIGN.CENTER)


def build_slide_6(prs):
    slide = add_content_slide(prs)
    set_title(slide, "보상 — 골드 & 명성")

    add_textbox(
        slide,
        Inches(0.6), Inches(1.15), Inches(12.0), Inches(0.45),
        "의뢰 완수 시 골드와 명성을 받습니다.",
        size=15, color=MUTED, align=PP_ALIGN.CENTER,
    )

    col_w = Inches(5.6)
    left_x = Inches(0.7)
    right_x = Inches(6.6)
    top = Inches(1.75)
    h = Inches(4.5)

    for x, title, subtitle, items, color in [
        (left_x, "골드", "수평 컨텐츠", ["공장 구역 해금", "기초자원 긴급 구매", "공장 확장·운영"], GOLD),
        (right_x, "명성", "수직 컨텐츠", ["새 기계·레시피 해금", "더 어려운 의뢰 등장", "그만큼 높은 보수"], REP),
    ]:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        add_textbox(slide, x, top + Inches(0.25), col_w, Inches(0.45), title, size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(0.75), col_w, Inches(0.35), subtitle, size=14, color=MUTED, align=PP_ALIGN.CENTER)

        tf_box = slide.shapes.add_textbox(x + Inches(0.35), top + Inches(1.25), col_w - Inches(0.7), Inches(2.8))
        tf = tf_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = f"• {item}"
            style_run(run, size=15)


def build_slide_7(prs):
    slide = add_content_slide(prs)
    set_title(slide, "스토리 & 목표")

    bullets = [
        "의뢰 중 반드시 완수해야 하는 스토리 의뢰가 있습니다.",
        "스토리 의뢰 실패 → 게임 오버 → 해당 일차부터 재도전",
        "최종 목표: 모든 스토리 의뢰 클리어",
    ]
    tf_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.0))
    tf = tf_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = f"• {line}"
        style_run(run, size=16)

    flow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.5), Inches(2.8))
    flow.fill.solid()
    flow.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFA)
    flow.line.color.rgb = ACCENT

    add_textbox(
        slide,
        Inches(1.0), Inches(3.75), Inches(11.1), Inches(2.3),
        "공장 운영  →  의뢰 해결  →  보상 획득  →  공장 성장\n"
        "        ↓\n"
        "더 어려운 의뢰  →  스토리 의뢰 물품 생산  →  엔딩",
        size=17, bold=False, color=DARK, align=PP_ALIGN.CENTER,
    )


def build_slide_8(prs):
    slide = add_content_slide(prs)
    set_title(slide, "팀 모집 — 직군별 역할")

    add_textbox(
        slide,
        Inches(0.6), Inches(1.15), Inches(12.0), Inches(0.4),
        "기획 · 개발 · 아트 각 1명씩 모집합니다.",
        size=15, color=MUTED, align=PP_ALIGN.CENTER,
    )

    rows, cols = 4, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.65), Inches(11.7), Inches(4.2)).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(9.7)

    for i, h in enumerate(["직군", "담당 업무"]):
        fill_table_cell(table.cell(0, i), h, size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    shade_header_row(table, cols)

    roles = [
        ("기획", "상세 레벨 디자인  (컨텐츠는 거의 구상 완료)"),
        ("개발", "의뢰·납품 시스템 · 세이브 구조 · UI  (공장 바깥)"),
        ("아트", "주인공 모션 · 주요 인물 초상화  (기타는 에셋·AI 활용)"),
    ]
    for r, (role, desc) in enumerate(roles, start=1):
        fill_table_cell(table.cell(r, 0), role, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        fill_table_cell(table.cell(r, 1), desc, size=14)


def build_slide_9(prs):
    slide = add_content_slide(prs)
    set_title(slide, "목표 & 확장 컨텐츠")

    add_textbox(
        slide,
        Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.85),
        "🎯 UNICON까지 완성 출품  |  MVP: 뒷산 동굴 (기믹 없는 공장 1개)\n"
        "지금은 좁은 범위에 집중 — 확장 시 기믹 던전·공장 간 시너지 예정",
        size=14, color=MUTED,
    )

    rows, cols = 7, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.3)).table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(7.3)

    headers = ["던전", "역할", "핵심 기믹"]
    for i, h in enumerate(headers):
        fill_table_cell(table.cell(0, i), h, size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    shade_header_row(table, cols)

    dungeons = [
        ("뒷산 동굴", "MVP·튜토리얼", "기믹 없음"),
        ("혹한의 설원", "환경 제약", "냉기·결빙·수동 해빙"),
        ("슬라임 서식지", "단일 원료", "슬라임 1종 → 다품종 가공"),
        ("깊이 잠든 황금향", "경제 특수", "명성 없음, 골드로 해결"),
        ("좋은 대장간?", "의뢰 해석", "애매한·재량 의뢰"),
        ("태고의 땅", "생물 공정", "포획·사육·생물 투입"),
    ]
    for r, row in enumerate(dungeons, start=1):
        bold_first = r == 1
        for c, val in enumerate(row):
            fill_table_cell(
                table.cell(r, c), val, size=12,
                bold=bold_first and c == 0,
                color=ACCENT if bold_first and c == 0 else DARK,
            )
        if r == 1:
            for c in range(3):
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xF8)


def build_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    if slide.shapes.title:
        slide.shapes.title.text = "DungeonFront"
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = (
                "기획 발표를 들어주셔서 감사합니다.\n"
                "재미있는 게임을 완성할 수 있도록 잘 리드하겠습니다.\n\n"
                "기획 · 개발 · 아트 — 많이 지원해 주세요!"
            )
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    style_run(run, size=20, color=MUTED)


def main():
    prs = Presentation(SRC)

    # 슬라이드 4~5(기존 초안) 제거 후 4~10 신규 생성
    while len(prs.slides) > 3:
        delete_slide(prs, len(prs.slides) - 1)

    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)

    prs.save(DST)
    print(f"Done: {len(prs.slides)} slides saved to {DST}")


if __name__ == "__main__":
    main()
